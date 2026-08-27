"""
Lector de presentaciones PowerPoint (Taller Práctico - Parte 1)

Extrae texto de archivos .pptx para poder indexarlos y buscar
información dentro de la presentación por similitud semántica.
"""
from typing import List, Dict


class PresentationReader:
    """Lee y extrae texto de archivos PowerPoint (.pptx)"""

    def __init__(self, pptx_path: str):
        """
        Args:
            pptx_path: Ruta al archivo .pptx
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx no está instalado. Instálalo con: pip install python-pptx"
            )

        self.pptx_path = pptx_path
        self.presentation = Presentation(pptx_path)

    def _extract_text_from_shape(self, shape) -> str:
        """Extrae el texto de una forma (incluye tablas)"""
        texts = []

        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = "".join(run.text for run in paragraph.runs)
                if paragraph_text.strip():
                    texts.append(paragraph_text.strip())

        if shape.has_table:
            for row in shape.table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    texts.append(" | ".join(row_texts))

        return "\n".join(texts)

    def read_slide(self, slide_number: int) -> Dict:
        """
        Lee el contenido de una diapositiva específica

        Returns:
            Diccionario con número, título y contenido de la diapositiva
        """
        slide = self.presentation.slides[slide_number - 1]

        title = ""
        content = []

        for shape in slide.shapes:
            shape_text = self._extract_text_from_shape(shape)

            if not shape_text:
                continue

            # El primer texto que parece título se usa como título
            if not title and len(shape_text) < 80:
                title = shape_text
            else:
                content.append(shape_text)

        return {
            "slide_number": slide_number,
            "title": title,
            "content": content
        }

    def read_all_slides(self) -> List[Dict]:
        """Lee todas las diapositivas de la presentación"""
        slides_data = []

        for i in range(1, len(self.presentation.slides) + 1):
            try:
                slides_data.append(self.read_slide(i))
            except Exception as e:
                print(f"⚠️  Error leyendo diapositiva {i}: {e}")

        return slides_data


def main():
    """Función de prueba"""
    import sys

    print("=" * 70)
    print("LECTOR DE PRESENTACIONES - PRUEBA")
    print("=" * 70 + "\n")

    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    else:
        pptx_path = "../Sesion6_Recuperación_de_información.pptx"

    try:
        reader = PresentationReader(pptx_path)
        slides = reader.read_all_slides()

        print(f"✓ {len(slides)} diapositivas leídas\n")

        for slide in slides[:5]:  # Mostrar las primeras 5
            print(f"Diapositiva {slide['slide_number']}: {slide['title']}")
            if slide['content']:
                print(f"   {slide['content'][0][:100]}...")
            print()

    except FileNotFoundError:
        print(f"❌ No se encontró: {pptx_path}")
    except Exception as e:
        print(f"❌ Error: {e}")


if __name__ == "__main__":
    main()
